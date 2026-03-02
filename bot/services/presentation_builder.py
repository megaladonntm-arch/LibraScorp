from __future__ import annotations

import asyncio
import random
import re
import tempfile
from datetime import datetime
from pathlib import Path

from PIL import Image, ImageStat
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Pt

from bot.services.ai_text_presentation_generator import (
    SlideContent,
    resolve_pdf_template_asset,
    resolve_template_asset,
)

try:
    import fitz  # type: ignore
except Exception:  # pragma: no cover
    fitz = None

TITLE_ZONE = (0.08, 0.08, 0.84, 0.16)
FIRST_SLIDE_TITLE_ZONE = (0.08, 0.03, 0.84, 0.12)
BODY_ZONE = (0.10, 0.26, 0.80, 0.58)
BODY_ZONE_NO_TITLE = (0.08, 0.10, 0.84, 0.78)
TITLE_ZONE_CANDIDATES = (
    (0.07, 0.06, 0.86, 0.18),
    (0.08, 0.08, 0.84, 0.16),
    (0.10, 0.09, 0.80, 0.16),
    (0.12, 0.10, 0.76, 0.16),
    (0.10, 0.14, 0.80, 0.16),
)
BODY_ZONE_CANDIDATES = (
    (0.08, 0.24, 0.84, 0.60),
    (0.10, 0.26, 0.80, 0.58),
    (0.10, 0.30, 0.80, 0.54),
    (0.12, 0.28, 0.76, 0.56),
    (0.12, 0.34, 0.76, 0.50),
    (0.08, 0.32, 0.84, 0.52),
)

# 13 visual variants; if slides >13, style repeats by modulo.
THEMES = [
    ("#F5F7FF", "#2F4B7C", "#6EA8FE", "#FFFFFF"),
    ("#FFF6EA", "#8C4A1A", "#F4A261", "#FFFFFF"),
    ("#EEF9F1", "#1B6B4A", "#7BCFA3", "#FFFFFF"),
    ("#FFF0F4", "#8E204B", "#E87EA1", "#FFFFFF"),
    ("#F2F3F7", "#30343F", "#8D99AE", "#FFFFFF"),
    ("#F9F4FF", "#4C2A85", "#A98ED6", "#FFFFFF"),
    ("#ECFBFF", "#0F5B6E", "#5EC9E2", "#FFFFFF"),
    ("#FFF8E7", "#7B5A12", "#DDBB5A", "#FFFFFF"),
    ("#F1FFF8", "#215E46", "#6ED7A7", "#FFFFFF"),
    ("#FFF1EC", "#7A2E1C", "#E78A70", "#FFFFFF"),
    ("#F0F5FF", "#1E3A8A", "#7FA8FF", "#FFFFFF"),
    ("#F7FFF3", "#355E1D", "#9FCF5B", "#FFFFFF"),
    ("#FFF2FA", "#6C1D45", "#D77AB3", "#FFFFFF"),
]

PDF_RENDER_SCALE = 1.8


def _safe_filename(source: str) -> str:
    cleaned = re.sub(r"[^\w\-]+", "_", source, flags=re.UNICODE).strip("_")
    return cleaned[:40] or "presentation"


def _parse_hex_color(color_hex: str) -> RGBColor:
    value = color_hex.strip().lstrip("#")
    if len(value) != 6:
        raise ValueError("Некорректный цвет")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _rgb_from_hex(color_hex: str) -> RGBColor:
    value = color_hex.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _estimate_body_font_size(slide: SlideContent) -> int:
    max_len = max((len(item) for item in slide.bullets), default=0)
    total_len = sum(len(item) for item in slide.bullets)
    bullet_count = len(slide.bullets)
    if bullet_count >= 5 or max_len > 210 or total_len > 680:
        return 16
    if bullet_count >= 4 or max_len > 170 or total_len > 520:
        return 18
    return 19


def _estimate_title_font_size(title: str) -> int:
    length = len(title.strip())
    if length > 90:
        return 24
    if length > 70:
        return 26
    if length > 52:
        return 30
    return 33


def _fit_inside(width: int, height: int, max_width: int, max_height: int) -> tuple[int, int]:
    if width <= 0 or height <= 0 or max_width <= 0 or max_height <= 0:
        return max(1, max_width), max(1, max_height)
    ratio = min(max_width / float(width), max_height / float(height))
    return max(1, int(width * ratio)), max(1, int(height * ratio))


def _add_background_image_cover(
    slide,
    image_path: Path,
    slide_width: int,
    slide_height: int,
) -> None:
    # Render full-bleed background while preserving original proportions.
    picture = slide.shapes.add_picture(
        str(image_path),
        left=0,
        top=0,
        width=slide_width,
        height=slide_height,
    )
    try:
        with Image.open(image_path) as img:
            img_width, img_height = img.size
    except Exception:
        return

    if img_width <= 0 or img_height <= 0 or slide_width <= 0 or slide_height <= 0:
        return

    img_ratio = img_width / float(img_height)
    slide_ratio = slide_width / float(slide_height)

    if img_ratio > slide_ratio:
        visible = slide_ratio / img_ratio
        crop_each = max(0.0, min(0.49, (1.0 - visible) / 2.0))
        picture.crop_left = crop_each
        picture.crop_right = crop_each
    elif img_ratio < slide_ratio:
        visible = img_ratio / slide_ratio
        crop_each = max(0.0, min(0.49, (1.0 - visible) / 2.0))
        picture.crop_top = crop_each
        picture.crop_bottom = crop_each


def _pick_image_layout(image_path: Path) -> str:
    layouts = ("left", "right", "top", "bottom")
    rng = random.SystemRandom()
    try:
        with Image.open(image_path) as img:
            img_width, img_height = img.size
    except Exception:
        return rng.choice(layouts)

    if img_width <= 0 or img_height <= 0:
        return rng.choice(layouts)

    # Keep random feel, but bias by image aspect for better visual balance.
    if img_width >= img_height:
        weighted = ("top", "bottom", "right", "left", "top")
    else:
        weighted = ("left", "right", "top", "bottom", "right")
    return rng.choice(weighted)


def _select_slide_image_pair(
    prepared_user_images: list[Path],
    slide_index: int,
) -> tuple[Path | None, Path | None]:
    if not prepared_user_images:
        return None, None
    if len(prepared_user_images) == 1:
        image = prepared_user_images[0]
        return image, image
    first_idx = (slide_index * 2) % len(prepared_user_images)
    second_idx = (first_idx + 1) % len(prepared_user_images)
    return prepared_user_images[first_idx], prepared_user_images[second_idx]


def _adjust_zones_for_dual_images(
    has_title: bool,
    preferred_large_layout: str,
) -> tuple[
    tuple[float, float, float, float],
    tuple[float, float, float, float],
    tuple[float, float, float, float],
]:
    if preferred_large_layout == "top":
        text_zone = (0.08, 0.50 if has_title else 0.44, 0.84, 0.42 if has_title else 0.48)
        large_zone = (0.06, 0.10 if has_title else 0.06, 0.88, 0.30 if has_title else 0.34)
        small_zone = (0.68, 0.38 if has_title else 0.34, 0.24, 0.16)
        return text_zone, large_zone, small_zone

    if preferred_large_layout == "bottom":
        text_zone = (0.08, 0.14 if has_title else 0.10, 0.84, 0.44 if has_title else 0.48)
        large_zone = (0.06, 0.62, 0.88, 0.30)
        small_zone = (0.08, 0.48, 0.22, 0.14)
        return text_zone, large_zone, small_zone

    if has_title:
        text_top = 0.24
        text_height = 0.64
        large_top = 0.18
        large_height = 0.70
    else:
        text_top = 0.10
        text_height = 0.78
        large_top = 0.08
        large_height = 0.82

    if preferred_large_layout == "left":
        text_zone = (0.53, text_top, 0.40, text_height)
        large_zone = (0.05, large_top, 0.45, large_height)
        small_zone = (0.34, 0.62 if has_title else 0.66, 0.16, 0.22)
    else:
        text_zone = (0.07, text_top, 0.40, text_height)
        large_zone = (0.50, large_top, 0.45, large_height)
        small_zone = (0.50, 0.62 if has_title else 0.66, 0.16, 0.22)

    return text_zone, large_zone, small_zone


def _add_user_image(
    slide,
    image_path: Path,
    image_zone: tuple[float, float, float, float],
    slide_width: int,
    slide_height: int,
) -> None:
    if not image_path.exists():
        return
    zone_left, zone_top, zone_width, zone_height = _ratio_to_emu(image_zone, slide_width, slide_height)
    if zone_width <= 0 or zone_height <= 0:
        return
    try:
        with Image.open(image_path) as img:
            image_width, image_height = img.size
    except Exception:
        return

    picture = slide.shapes.add_picture(
        str(image_path),
        left=zone_left,
        top=zone_top,
        width=zone_width,
        height=zone_height,
    )
    zone_ratio = zone_width / float(zone_height)
    image_ratio = image_width / float(image_height)
    if image_ratio > zone_ratio:
        visible = zone_ratio / image_ratio
        crop_each = max(0.0, min(0.49, (1.0 - visible) / 2.0))
        picture.crop_left = crop_each
        picture.crop_right = crop_each
    elif image_ratio < zone_ratio:
        visible = image_ratio / zone_ratio
        crop_each = max(0.0, min(0.49, (1.0 - visible) / 2.0))
        picture.crop_top = crop_each
        picture.crop_bottom = crop_each


def _ratio_to_emu(
    box: tuple[float, float, float, float],
    slide_width: int,
    slide_height: int,
) -> tuple[int, int, int, int]:
    left = int(slide_width * box[0])
    top = int(slide_height * box[1])
    width = int(slide_width * box[2])
    height = int(slide_height * box[3])
    return left, top, width, height


def _hex_to_rgb(color_hex: str) -> tuple[int, int, int]:
    value = color_hex.strip().lstrip("#")
    if len(value) != 6:
        return 0, 0, 0
    return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)


def _contrast_ratio(text_rgb: tuple[int, int, int], bg_luma: float) -> float:
    text_luma = (0.2126 * text_rgb[0] + 0.7152 * text_rgb[1] + 0.0722 * text_rgb[2]) / 255.0
    bg_norm = bg_luma / 255.0
    lighter = max(text_luma, bg_norm) + 0.05
    darker = min(text_luma, bg_norm) + 0.05
    return lighter / darker


def _score_candidate(image: Image.Image, box: tuple[float, float, float, float], text_rgb: tuple[int, int, int]) -> float:
    width, height = image.size
    left = int(width * box[0])
    top = int(height * box[1])
    right = int(width * (box[0] + box[2]))
    bottom = int(height * (box[1] + box[3]))

    if right <= left or bottom <= top:
        return float("-inf")

    cropped = image.crop((left, top, right, bottom)).convert("L")
    stat = ImageStat.Stat(cropped)
    mean = stat.mean[0] if stat.mean else 128.0
    stddev = stat.stddev[0] if stat.stddev else 0.0
    contrast = _contrast_ratio(text_rgb, mean)

    # Prefer areas with strong text contrast and lower visual noise.
    return (contrast * 120.0) + (255.0 - stddev) + (box[2] * box[3] * 25.0)


def _detect_text_zones_from_background(
    image_path: Path,
    text_color_hex: str,
) -> tuple[tuple[float, float, float, float], tuple[float, float, float, float]]:
    try:
        with Image.open(image_path) as image:
            rgb = _hex_to_rgb(text_color_hex)
            best_title = max(TITLE_ZONE_CANDIDATES, key=lambda candidate: _score_candidate(image, candidate, rgb))
            body_candidates = [candidate for candidate in BODY_ZONE_CANDIDATES if candidate[1] >= best_title[1] + best_title[3] - 0.01]
            if not body_candidates:
                body_candidates = list(BODY_ZONE_CANDIDATES)
            best_body = max(body_candidates, key=lambda candidate: _score_candidate(image, candidate, rgb))
            return best_title, best_body
    except Exception:
        return TITLE_ZONE, BODY_ZONE


def _theme_for_index(index: int) -> tuple[RGBColor, RGBColor, RGBColor, RGBColor]:
    bg_hex, header_hex, accent_hex, card_hex = THEMES[index % len(THEMES)]
    return (
        _rgb_from_hex(bg_hex),
        _rgb_from_hex(header_hex),
        _rgb_from_hex(accent_hex),
        _rgb_from_hex(card_hex),
    )


def _add_background(slide, slide_width: int, slide_height: int, index: int) -> tuple[RGBColor, RGBColor]:
    bg_color, header_color, accent_color, card_color = _theme_for_index(index)

    background = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left=0,
        top=0,
        width=slide_width,
        height=slide_height,
    )
    background.fill.solid()
    background.fill.fore_color.rgb = bg_color
    background.line.fill.background()

    top_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left=0,
        top=0,
        width=slide_width,
        height=int(slide_height * 0.10),
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = header_color
    top_bar.line.fill.background()

    decorative_circle = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        left=int(slide_width * 0.78),
        top=int(slide_height * 0.72),
        width=int(slide_width * 0.26),
        height=int(slide_width * 0.26),
    )
    decorative_circle.fill.solid()
    decorative_circle.fill.fore_color.rgb = accent_color
    decorative_circle.fill.transparency = 0.60
    decorative_circle.line.fill.background()

    content_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left=int(slide_width * 0.06),
        top=int(slide_height * 0.18),
        width=int(slide_width * 0.88),
        height=int(slide_height * 0.72),
    )
    content_card.fill.solid()
    content_card.fill.fore_color.rgb = card_color
    content_card.line.fill.solid()
    content_card.line.fill.fore_color.rgb = accent_color
    content_card.line.width = Pt(1.6)

    return header_color, accent_color


def _render_pdf_pages_to_png(pdf_path: Path) -> list[Path]:
    if fitz is None:
        raise RuntimeError("Для PDF-шаблона нужен пакет PyMuPDF (pip install pymupdf).")

    output: list[Path] = []
    document = fitz.open(str(pdf_path))
    matrix = fitz.Matrix(PDF_RENDER_SCALE, PDF_RENDER_SCALE)
    try:
        # For ready-made PDF templates, skip the first page/background consistently.
        start_page = 1 if document.page_count > 1 else 0
        for page_index in range(start_page, document.page_count):
            page = document.load_page(page_index)
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            image_path = Path(tempfile.mkstemp(prefix="pdf_template_", suffix=".png")[1])
            pixmap.save(str(image_path))
            output.append(image_path)
    finally:
        document.close()
    return output


def _build_presentation_sync(
    topic: str,
    template_types: list[int],
    slides: list[SlideContent],
    font_name: str,
    font_color: str,
    creator_names: str | None = None,
    creator_title: str = "Presentation creators",
    user_image_paths: list[str] | None = None,
) -> Path:
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    color = _parse_hex_color(font_color)
    temp_images: list[Path] = []
    pdf_pages_cache: dict[str, list[Path]] = {}
    zones_cache: dict[str, tuple[tuple[float, float, float, float], tuple[float, float, float, float]]] = {}
    prepared_user_images: list[Path] = []
    for image_path in user_image_paths or []:
        candidate = Path(image_path)
        if candidate.exists():
            prepared_user_images.append(candidate)

    for index, slide_content in enumerate(slides):
        template_type = template_types[index] if index < len(template_types) else (template_types[0] if template_types else 1)
        slide = presentation.slides.add_slide(blank_layout)

        pdf_template_path = resolve_pdf_template_asset(template_type)
        static_image_asset = resolve_template_asset(template_type)

        if pdf_template_path is not None:
            cache_key = str(pdf_template_path.resolve())
            pdf_pages = pdf_pages_cache.get(cache_key)
            if pdf_pages is None:
                pdf_pages = _render_pdf_pages_to_png(pdf_template_path)
                pdf_pages_cache[cache_key] = pdf_pages
                temp_images.extend(pdf_pages)
            if not pdf_pages:
                raise RuntimeError(f"PDF шаблон пустой: {pdf_template_path}")
            bg_image = pdf_pages[index % len(pdf_pages)]
            _add_background_image_cover(
                slide=slide,
                image_path=bg_image,
                slide_width=presentation.slide_width,
                slide_height=presentation.slide_height,
            )
            zone_key = str(bg_image.resolve())
        elif static_image_asset is not None:
            _add_background_image_cover(
                slide=slide,
                image_path=static_image_asset,
                slide_width=presentation.slide_width,
                slide_height=presentation.slide_height,
            )
            zone_key = str(static_image_asset.resolve())
        else:
            _add_background(slide, presentation.slide_width, presentation.slide_height, index)
            zone_key = ""

        if zone_key:
            zones = zones_cache.get(zone_key)
            if zones is None:
                background_path = Path(zone_key)
                zones = _detect_text_zones_from_background(background_path, font_color)
                zones_cache[zone_key] = zones
            title_zone, body_zone = zones
        else:
            title_zone, body_zone = TITLE_ZONE, BODY_ZONE

        has_title = index == 0
        title_text = topic if has_title else ""
        if has_title:
            title_zone = FIRST_SLIDE_TITLE_ZONE
        first_image, second_image = _select_slide_image_pair(prepared_user_images, index)
        if first_image is None and zone_key:
            fallback_image = Path(zone_key)
            if fallback_image.exists():
                first_image = fallback_image
                second_image = fallback_image
        body_zone_for_text = body_zone if has_title else BODY_ZONE_NO_TITLE
        large_image_zone: tuple[float, float, float, float] | None = None
        small_image_zone: tuple[float, float, float, float] | None = None
        if first_image is not None and second_image is not None:
            preferred_layout = _pick_image_layout(first_image)
            body_zone_for_text, large_image_zone, small_image_zone = _adjust_zones_for_dual_images(
                has_title=has_title,
                preferred_large_layout=preferred_layout,
            )

        if has_title:
            title_left, title_top, title_width, title_height = _ratio_to_emu(
                title_zone,
                presentation.slide_width,
                presentation.slide_height,
            )
            title_box = slide.shapes.add_textbox(
                left=title_left,
                top=title_top,
                width=title_width,
                height=title_height,
            )
            title_frame = title_box.text_frame
            title_frame.clear()
            title_frame.margin_left = 0
            title_frame.margin_right = 0
            title_frame.margin_top = 0
            title_frame.margin_bottom = 0
            title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            title_frame.word_wrap = True
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.text = title_text
            title_paragraph.font.bold = True
            title_paragraph.font.name = font_name
            title_paragraph.font.size = Pt(_estimate_title_font_size(title_text))
            title_paragraph.font.color.rgb = color
            title_paragraph.alignment = PP_ALIGN.CENTER
            title_paragraph.space_after = Pt(2)

        body_left, body_top, body_width, body_height = _ratio_to_emu(
            body_zone_for_text,
            presentation.slide_width,
            presentation.slide_height,
        )
        body_box = slide.shapes.add_textbox(
            left=body_left,
            top=body_top,
            width=body_width,
            height=body_height,
        )
        body_frame = body_box.text_frame
        body_frame.clear()
        body_frame.margin_left = 0
        body_frame.margin_right = 0
        body_frame.margin_top = 0
        body_frame.margin_bottom = 0
        body_frame.vertical_anchor = MSO_ANCHOR.TOP
        body_frame.word_wrap = True
        body_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        body_font_size = _estimate_body_font_size(slide_content)

        for bullet_index, bullet in enumerate(slide_content.bullets):
            paragraph = body_frame.paragraphs[0] if bullet_index == 0 else body_frame.add_paragraph()
            paragraph.text = f"• {bullet}"
            paragraph.level = 0
            paragraph.font.name = font_name
            paragraph.font.size = Pt(body_font_size)
            paragraph.font.color.rgb = color
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(5)

        if first_image is not None and large_image_zone is not None:
            _add_user_image(
                slide=slide,
                image_path=first_image,
                image_zone=large_image_zone,
                slide_width=presentation.slide_width,
                slide_height=presentation.slide_height,
            )
        if second_image is not None and small_image_zone is not None:
            _add_user_image(
                slide=slide,
                image_path=second_image,
                image_zone=small_image_zone,
                slide_width=presentation.slide_width,
                slide_height=presentation.slide_height,
            )

    if creator_names:
        creator_index = len(slides)
        final_slide = presentation.slides.add_slide(blank_layout)
        _add_background(
            final_slide,
            presentation.slide_width,
            presentation.slide_height,
            creator_index,
        )

        names_box = final_slide.shapes.add_textbox(
            left=int(presentation.slide_width * 0.10),
            top=int(presentation.slide_height * 0.34),
            width=int(presentation.slide_width * 0.80),
            height=int(presentation.slide_height * 0.36),
        )
        names_frame = names_box.text_frame
        names_frame.clear()
        names_frame.vertical_anchor = MSO_ANCHOR.TOP
        names_frame.word_wrap = True
        for idx, raw_name in enumerate([x.strip() for x in creator_names.split(",") if x.strip()][:20]):
            paragraph = names_frame.paragraphs[0] if idx == 0 else names_frame.add_paragraph()
            paragraph.text = raw_name
            paragraph.level = 0
            paragraph.font.name = font_name
            paragraph.font.size = Pt(26)
            paragraph.font.color.rgb = color
            paragraph.alignment = PP_ALIGN.CENTER

    out_dir = Path(tempfile.mkdtemp(prefix="tg_presentation_"))
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = out_dir / f"{_safe_filename(topic)}_{stamp}.pptx"
    presentation.save(output_path)
    for temp_image in temp_images:
        try:
            temp_image.unlink(missing_ok=True)
        except Exception:
            pass
    return output_path


async def build_presentation_file(
    topic: str,
    template_types: list[int],
    slides: list[SlideContent],
    font_name: str,
    font_color: str,
    creator_names: str | None = None,
    creator_title: str = "Presentation creators",
    user_image_paths: list[str] | None = None,
) -> Path:
    return await asyncio.to_thread(
        _build_presentation_sync,
        topic,
        template_types,
        slides,
        font_name,
        font_color,
        creator_names,
        creator_title,
        user_image_paths,
    )

